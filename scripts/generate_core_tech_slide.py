from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "StoryForge核心技术-AI_DM剧本生成与安全约束机制.pptx"

FONT = "Microsoft YaHei"
TITLE_FONT = "Microsoft YaHei UI"

BG = RGBColor(16, 12, 8)
PANEL = RGBColor(34, 26, 17)
PANEL_DARK = RGBColor(22, 17, 12)
GOLD = RGBColor(222, 169, 64)
GOLD_LIGHT = RGBColor(248, 212, 128)
TEXT = RGBColor(239, 229, 205)
MUTED = RGBColor(178, 155, 118)
WARN = RGBColor(116, 61, 31)


def add_textbox(slide, x, y, w, h, text, size=12, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, title, bullets, title_size=11, bullet_size=8.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_DARK
    shape.line.color.rgb = RGBColor(89, 65, 31)
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.06)

    p = frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = GOLD_LIGHT

    for item in bullets:
        p = frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(1)
    return shape


def add_card(slide, x, y, w, h, title, desc, accent=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(12.5)
    r.font.bold = True
    r.font.color.rgb = GOLD_LIGHT
    p = frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = desc
    r.font.name = FONT
    r.font.size = Pt(8.5)
    r.font.color.rgb = MUTED
    return shape


def add_flow_box(slide, x, y, w, h, title, desc="", highlight=False):
    fill = WARN if highlight else PANEL
    line = GOLD_LIGHT if highlight else GOLD
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.4)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(8.8)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if desc:
        p = frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = desc
        r.font.name = FONT
        r.font.size = Pt(6.8)
        r.font.color.rgb = GOLD_LIGHT if highlight else MUTED
    return shape


def add_arrow(slide, x, y, w, h, direction="right"):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }[direction]
    arrow = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.color.rgb = GOLD
    return arrow


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    # Outer archive frame
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(0.16), Inches(12.96), Inches(7.16))
    frame.fill.background()
    frame.line.color.rgb = RGBColor(57, 40, 20)
    frame.line.width = Pt(1)

    add_textbox(slide, 0.48, 0.35, 0.25, 0.35, "◆", size=22, color=GOLD_LIGHT, bold=True)
    add_textbox(
        slide,
        0.78,
        0.28,
        8.6,
        0.45,
        "核心技术：AI DM 剧本生成与安全约束机制",
        size=25,
        color=GOLD_LIGHT,
        bold=True,
    )
    add_textbox(
        slide,
        0.82,
        0.78,
        6.8,
        0.28,
        "AI API 不是简单聊天，而是嵌入游戏主流程的剧情推进引擎",
        size=10.5,
        color=MUTED,
    )

    # Left: input routing
    add_panel(slide, 0.55, 1.23, 2.35, 4.86, "输入分流机制", [], title_size=13)
    add_card(slide, 0.76, 1.82, 1.92, 0.74, "行动 DM", "推动剧情 / 进入判定")
    add_card(slide, 0.76, 2.84, 1.92, 0.74, "询问 DM", "非剧透提示 / 不改状态")
    add_card(slide, 0.76, 3.86, 1.92, 0.74, "普通聊天", "同步消息 / 不推进剧情")
    add_textbox(
        slide,
        0.72,
        5.04,
        1.98,
        0.58,
        "行动会改变剧情状态；询问只辅助理解；聊天只同步给房间成员。",
        size=8.3,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    # Center: main flow
    add_panel(slide, 3.08, 1.23, 6.12, 4.86, "行动 DM 主流程", [], title_size=13)
    xs = [3.32, 4.78, 6.24, 7.70]
    y1, y2 = 1.88, 3.02
    w, h = 1.22, 0.54
    add_flow_box(slide, xs[0], y1, w, h, "玩家行动", "submit")
    add_flow_box(slide, xs[1], y1, w, h, "ActionParser", "意图解析")
    add_flow_box(slide, xs[2], y1, w, h, "技能 / DC", "风险等级")
    add_flow_box(slide, xs[3], y1, w, h, "d20 判定", "rule_service", highlight=True)
    for ax in [4.53, 5.99, 7.45]:
        add_arrow(slide, ax, y1 + 0.18, 0.18, 0.18, "right")
    add_arrow(slide, xs[3] + 0.48, 2.48, 0.26, 0.34, "down")

    add_flow_box(slide, xs[3], y2, w, h, "NarrativeAgent", "叙事反馈")
    add_flow_box(slide, xs[2], y2, w, h, "CriticAgent", "合理性审核")
    add_flow_box(slide, xs[1], y2, w, h, "RevisionLoop", "必要时修正")
    add_flow_box(slide, xs[0], y2, w, h, "DB + WS", "存储并同步")
    for ax in [7.45, 5.99, 4.53]:
        add_arrow(slide, ax, y2 + 0.18, 0.18, 0.18, "left")

    rule = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.42), Inches(3.92), Inches(5.45), Inches(0.55))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(47, 33, 15)
    rule.line.color.rgb = GOLD_LIGHT
    rule.line.width = Pt(1.2)
    tf = rule.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "关键边界：AI 不直接决定成败；成败由后端 d20 判定，AI 负责叙事表达"
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = GOLD_LIGHT

    add_textbox(
        slide,
        3.34,
        4.72,
        5.56,
        0.54,
        "询问 DM：玩家提问 → GuidanceAgent → 读取房间/角色/线索 → 非剧透提示 → 不骰点 / 不改状态",
        size=8.5,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        3.52,
        5.38,
        5.2,
        0.36,
        "普通聊天：只进入 room_messages 与 WebSocket 广播，不进入剧情状态更新链路",
        size=7.8,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Right: constraints
    add_panel(
        slide,
        9.48,
        1.23,
        3.27,
        1.62,
        "Prompt 约束",
        [
            "遵守世界观、角色、历史事实",
            "不改写已发生剧情",
            "不绕过后端 d20 判定",
            "不直接给关键道具/通关/死亡",
            "输出反馈 + 下一步方向",
        ],
        title_size=12,
        bullet_size=7.5,
    )
    add_panel(
        slide,
        9.48,
        3.02,
        3.27,
        1.34,
        "保守机制",
        [
            "高风险行动必须判定",
            "主线剧情渐进式推进",
            "不确定信息用“似乎/可能”",
            "Critic 检查越界、冲突、剧透",
        ],
        title_size=12,
        bullet_size=7.6,
    )
    add_panel(
        slide,
        9.48,
        4.54,
        3.27,
        1.34,
        "AI 安全设计",
        [
            "防 Prompt 注入",
            "防越权操作",
            "防 Prompt / API Key / DB 泄露",
            "拒绝现实危险与违法指导",
        ],
        title_size=12,
        bullet_size=7.6,
    )

    # Bottom summary
    bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.38), Inches(12.2), Inches(0.62))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RGBColor(30, 21, 11)
    bottom.line.color.rgb = GOLD
    bottom.line.width = Pt(1.2)
    tf = bottom.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AI DM 的核心不是“会聊天”，而是把玩家输入转化为可解析、可判定、可存储、可同步、可约束的剧情推进流程。"
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = GOLD_LIGHT

    add_textbox(slide, 12.66, 7.05, 0.35, 0.18, "1", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
